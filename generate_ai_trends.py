#!/usr/bin/env python3
"""
最近のAI事情 - スライド生成スクリプト
pptx-generatorスキルに基づき python-pptx で作成
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ========== スライド1: タイトル ==========
title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "最近のAI事情 2025"
p.font.size = Pt(44)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

sub_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.text = "生成AI・マルチモーダル・AIエージェントの最新トレンド"
p2.font.size = Pt(24)
p2.alignment = PP_ALIGN.CENTER

# ========== スライド2: 2025年AI 5大トレンド ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "2025年 AI 5大トレンド"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
items = [
    ("1. AIエージェントの本格実装", "質問に答えるだけでなく、業務プロセスを自律的に遂行"),
    ("2. 軽量・効率化の劇的進展", "単一GPUで動作する軽量モデルが登場"),
    ("3. 現場実装と生産性の可視化", "導入効果が明確な数値で語られるように"),
    ("4. マルチモーダル・ロボティクス融合", "テキスト・画像・音声・動画を横断処理"),
    ("5. ガバナンスとガードレール強化", "倫理的ガバナンスの整備が進行中"),
]
for i, (title_text, desc) in enumerate(items):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"{title_text}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = f"  → {desc}"
    p2.font.size = Pt(14)
    p2.level = 1
    p2.space_after = Pt(12)

# ========== スライド3: AIエージェント ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "AIエージェントの本格実装"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True
points = [
    "OpenAI「Operator」、Salesforce「Agentforce 2dx」など",
    "ブラウザ操作・CRM統合による業務自動化が実現",
    "国内：NTTデータ「LITRON Marketing」がマーケティング業務を約6割削減",
    "単なるチャットボットから「自律的にタスクを完遂する」段階へ",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド4: マルチモーダルAI ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "マルチモーダルAIの台頭"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "テキスト・画像・音声・動画を横断的に処理・生成",
    "Sora（動画）、Midjourney（画像）、Suno（音楽）など",
    "人間のように「世界を総合的に解釈する」能力へ",
    "医療画像診断、教育、エンタメなど幅広い分野で活用",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド5: 主要モデル比較 ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "主要生成AIモデル比較（2025年）"
p.font.size = Pt(32)
p.font.bold = True

# テーブル作成
rows, cols = 4, 4
left = Inches(0.5)
top = Inches(1.3)
width = Inches(9)
height = Inches(4)
table = content_slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.5)
table.columns[3].width = Inches(2.5)

headers = ["モデル", "強み", "料金目安", "おすすめ用途"]
data = [
    ["GPT-5.2", "運用・自動化に最適\n推論速度18%向上", "約3,000円/月", "反復タスク・自動化"],
    ["Claude Opus 4.5", "文章品質・コーディング\n200kトークン対応", "約3,000円/月", "提案書・長文分析"],
    ["Gemini 3", "最高コスパ\n1Mトークン対応", "2,900円/月", "速度・大量処理・Workspace連携"],
]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(14)
for r, row in enumerate(data, 1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(12)

# ========== スライド6: 軽量化トレンド ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "軽量・効率化の劇的進展"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "Gemma 3、Sakana AI「TinySwallow-1.5B」など",
    "単一GPUで動作する軽量モデルが実用レベルに",
    "パナソニックHD：画像生成効率を5倍に高める技術を発表",
    "「高性能＝巨大モデル」という前提が崩れつつある",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド7: 社会実装の課題 ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "社会実装の課題とガバナンス"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "AI生成物の透明性確保、プライバシー保護、著作権対応",
    "感覚や主観の完全再現は現状困難",
    "データ偏在、倫理的考慮が重要な課題",
    "「単なる技術ツール」から「社会基盤」への転換が本格化",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド8: まとめ ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "まとめ"
p.font.size = Pt(36)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
tf = content_box.text_frame
p = tf.paragraphs[0]
p.text = "2025年は生成AIが「社会基盤」へ転換した年"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "• AIエージェントによる業務自動化が本格化"
p2.font.size = Pt(18)
p2.space_after = Pt(10)
p3 = tf.add_paragraph()
p3.text = "• マルチモーダルAIがクリエイティブ・医療・教育を変革"
p3.font.size = Pt(18)
p3.space_after = Pt(10)
p4 = tf.add_paragraph()
p4.text = "• 軽量化によりオンプレ・エッジでの活用が拡大"
p4.font.size = Pt(18)
p4.space_after = Pt(10)
p5 = tf.add_paragraph()
p5.text = "• 用途に応じたモデル選びが重要"
p5.font.size = Pt(18)

# 保存（output/ディレクトリへ）
output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "最近のAI事情_2025.pptx")
prs.save(output_path)
print(f"スライドを生成しました: {output_path}")
