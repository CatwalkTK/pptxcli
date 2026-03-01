"""PPTX AI Assistant CLI 仕様紹介プレゼンテーション生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== Color Palette ==========
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
SUBTLE_GRAY = RGBColor(0x64, 0x74, 0x8B)
TOTAL_SLIDES = 10


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=18, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        run_b = p.add_run()
        run_b.text = "\u25cf  "
        run_b.font.size = Pt(font_size - 4)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = 'Segoe UI'
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = 'Segoe UI'
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f'{num} / {total}', font_size=11, color=SUBTLE_GRAY,
             alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_shape(s, Inches(-1), Inches(-1), Inches(5), Inches(5), BG_CARD)
add_shape(s, Inches(10), Inches(4), Inches(5), Inches(5), BG_CARD)
add_accent_line(s, Inches(1.5), Inches(2.5), Inches(3), ACCENT_BLUE)
add_text(s, Inches(1.5), Inches(1.3), Inches(2), Inches(1),
         '>_', font_size=48, color=ACCENT_BLUE, bold=True, font_name='Consolas')
add_text(s, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
         'PPTX AI Assistant CLI', font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1.5), Inches(4.0), Inches(8), Inches(0.8),
         '自然言語でプレゼンテーションを生成する AI エージェント CLI ツール',
         font_size=22, color=LIGHT_GRAY)
for i, tag in enumerate(['Python', 'Claude AI', 'PowerPoint', 'Agent Loop']):
    x = Inches(1.5) + Inches(i * 2.2)
    add_shape(s, x, Inches(5.2), Inches(2.0), Inches(0.5), BG_CARD, ACCENT_BLUE)
    add_text(s, x, Inches(5.2), Inches(2.0), Inches(0.5),
             tag, font_size=13, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER, font_name='Consolas')
add_text(s, Inches(1.5), Inches(6.2), Inches(3), Inches(0.4),
         'v0.1.0  |  2026', font_size=14, color=SUBTLE_GRAY)
add_slide_number(s, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2: Overview
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
         '概要', font_size=36, color=WHITE, bold=True)
add_shape(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.8),
          BG_CARD, border_color=RGBColor(0x2A, 0x3A, 0x5C))
add_text(s, Inches(1.3), Inches(2.2), Inches(10.5), Inches(1.5),
         'Claude AI と PowerPoint 生成機能を統合した、インタラクティブな REPL ベースの\n'
         'エージェント CLI。自然言語の指示から高品質なプレゼンテーションを自動生成します。',
         font_size=20, color=LIGHT_GRAY)
features = [
    ('\U0001f916', 'AI エージェント',
     'Claude API による\n自律的なツール実行', ACCENT_BLUE),
    ('\U0001f4ca', 'PPTX 生成',
     '自然言語から\nプレゼン自動作成', ACCENT_GREEN),
    ('\U0001f527', '6種のツール',
     'ファイル操作・コマンド\n実行・PPTX生成', ACCENT_PURPLE),
    ('\U0001f4ac', 'REPL インターフェース',
     'Rich による美しい\nターミナル UI', ACCENT_ORANGE),
]
for i, (icon, title, desc, color) in enumerate(features):
    x = Inches(0.8) + Inches(i * 3.0)
    y = Inches(4.3)
    add_shape(s, x, y, Inches(2.8), Inches(2.5), BG_CARD, border_color=color)
    add_text(s, x, y + Inches(0.2), Inches(2.8), Inches(0.6),
             icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(0.9), Inches(2.4), Inches(0.5),
             title, font_size=17, color=color, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.4), Inches(2.4), Inches(0.9),
             desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(s, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3: Tech Stack
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
         '技術スタック', font_size=36, color=WHITE, bold=True)
techs = [
    ('Python 3.8+', '実行環境・コア言語', 'async/await, type hints', ACCENT_BLUE),
    ('Anthropic Claude API', 'AI エンジン', 'AsyncAnthropic, Tool Use', ACCENT_GREEN),
    ('python-pptx', 'PPTX 生成', 'OOXML ベースの直接生成', ACCENT_PURPLE),
    ('Rich', 'ターミナル UI', 'Panel, Markdown, Console', ACCENT_ORANGE),
    ('prompt_toolkit', '入力インターフェース', 'PromptSession, History', ACCENT_PINK),
    ('python-dotenv', '設定管理', '.env ファイルの読み込み', ACCENT_BLUE),
]
for i, (name, role, detail, color) in enumerate(techs):
    col, row = i % 3, i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.0) + Inches(row * 2.7)
    add_shape(s, x, y, Inches(3.8), Inches(2.3), BG_CARD, border_color=color)
    add_shape(s, x, y, Pt(4), Inches(2.3), color)
    add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.5),
             name, font_size=20, color=color, bold=True, font_name='Consolas')
    add_text(s, x + Inches(0.3), y + Inches(0.8), Inches(3.3), Inches(0.4),
             role, font_size=16, color=WHITE)
    add_text(s, x + Inches(0.3), y + Inches(1.3), Inches(3.3), Inches(0.8),
             detail, font_size=14, color=LIGHT_GRAY)
add_slide_number(s, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4: Architecture
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_PURPLE)
add_text(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
         'アーキテクチャ', font_size=36, color=WHITE, bold=True)
flow = [
    ('User Input', 'ユーザー入力', ACCENT_BLUE),
    ('REPL Loop', 'app.py', ACCENT_GREEN),
    ('Agent', 'agent.py', ACCENT_PURPLE),
    ('Claude API', 'Anthropic', ACCENT_ORANGE),
    ('Tool 実行', 'tools/', ACCENT_PINK),
    ('UI Render', 'ui/', ACCENT_BLUE),
]
yc = Inches(3.0)
for i, (label, sub, color) in enumerate(flow):
    x = Inches(0.5) + Inches(i * 2.15)
    add_shape(s, x, yc, Inches(1.9), Inches(1.4), BG_CARD, border_color=color)
    add_text(s, x, yc + Inches(0.15), Inches(1.9), Inches(0.5),
             label, font_size=16, color=color, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Consolas')
    add_text(s, x, yc + Inches(0.65), Inches(1.9), Inches(0.5),
             sub, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text(s, x + Inches(1.95), yc + Inches(0.3), Inches(0.3), Inches(0.5),
                 '\u2192', font_size=24, color=SUBTLE_GRAY,
                 alignment=PP_ALIGN.CENTER)
add_shape(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8),
          BG_CARD, border_color=ACCENT_PURPLE)
add_text(s, Inches(1.2), Inches(5.1), Inches(4), Inches(0.5),
         'エージェントループ（最大20回）',
         font_size=18, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(1.2), Inches(5.6), Inches(10.8), Inches(1.0),
         'Claude API レスポンスに tool_use が含まれる限り自動的にツールを実行し、'
         '結果を次のリクエストに含めて\n'
         '再送信するループ。stop_reason が end_turn になるまで繰り返し、'
         '最終テキストをユーザーに返却。',
         font_size=15, color=LIGHT_GRAY)
add_slide_number(s, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5: Tools (1/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_ORANGE)
add_text(s, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
         '搭載ツール  (1/2) — ファイル操作', font_size=36, color=WHITE, bold=True)
file_tools = [
    ('read_file', 'ファイル読み取り',
     ['path を指定してファイル内容を取得',
      '最大 100,000 文字まで読み込み',
      'テキストファイル全般に対応'], ACCENT_BLUE),
    ('write_file', 'ファイル作成・上書き',
     ['path と content を指定して書き込み',
      '親ディレクトリを自動作成（makedirs）',
      '文字数を含む確認メッセージを返却'], ACCENT_GREEN),
    ('list_files', 'ディレクトリ一覧',
     ['path（デフォルト: "."）のファイル一覧',
      'recursive オプションで再帰探索',
      '最大 500 件、隠しフォルダ除外'], ACCENT_PURPLE),
]
for i, (name, desc, details, color) in enumerate(file_tools):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.0)
    add_shape(s, x, y, Inches(3.8), Inches(5.0), BG_CARD, border_color=color)
    add_shape(s, x, y, Inches(3.8), Pt(4), color)
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.5),
             name, font_size=22, color=color, bold=True, font_name='Consolas')
    add_text(s, x + Inches(0.3), y + Inches(0.9), Inches(3.3), Inches(0.5),
             desc, font_size=17, color=WHITE, bold=True)
    add_bullet_text(s, x + Inches(0.3), y + Inches(1.6), Inches(3.3), Inches(3.0),
                    details, font_size=15, color=LIGHT_GRAY, bullet_color=color)
add_slide_number(s, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6: Tools (2/2)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_ORANGE)
add_text(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7),
         '搭載ツール  (2/2) — 検索・コマンド・PPTX',
         font_size=36, color=WHITE, bold=True)
more_tools = [
    ('search_files', 'Globパターン検索',
     ['pattern でファイルを検索',
      '例: **/*.py, *.md',
      '最大 200 件の結果を返却'], ACCENT_ORANGE),
    ('execute_command', 'シェルコマンド実行',
     ['command 文字列を実行',
      'タイムアウト 30秒（変更可）',
      '出力最大 50,000 文字'], ACCENT_PINK),
    ('generate_pptx', 'PPTX プレゼン生成',
     ['filename, title, slides を指定',
      'content / two_column 2タイプ',
      '4x3比率・自動フォーマット'], ACCENT_BLUE),
]
for i, (name, desc, details, color) in enumerate(more_tools):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.0)
    add_shape(s, x, y, Inches(3.8), Inches(5.0), BG_CARD, border_color=color)
    add_shape(s, x, y, Inches(3.8), Pt(4), color)
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.5),
             name, font_size=22, color=color, bold=True, font_name='Consolas')
    add_text(s, x + Inches(0.3), y + Inches(0.9), Inches(3.3), Inches(0.5),
             desc, font_size=17, color=WHITE, bold=True)
    add_bullet_text(s, x + Inches(0.3), y + Inches(1.6), Inches(3.3), Inches(3.0),
                    details, font_size=15, color=LIGHT_GRAY, bullet_color=color)
add_slide_number(s, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7: CLI Commands
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_PINK)
add_text(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
         'CLI コマンド & 操作', font_size=36, color=WHITE, bold=True)

# Terminal card
add_shape(s, Inches(0.8), Inches(2.0), Inches(7.5), Inches(5.0),
          RGBColor(0x11, 0x14, 0x22), border_color=SUBTLE_GRAY)
add_shape(s, Inches(0.8), Inches(2.0), Inches(7.5), Inches(0.45),
          RGBColor(0x1E, 0x29, 0x3B))
add_text(s, Inches(1.1), Inches(2.03), Inches(4), Inches(0.4),
         '\u25cf  \u25cf  \u25cf   PPTX AI Assistant',
         font_size=12, color=SUBTLE_GRAY)
cmds = [
    ('$ ', 'python -m cli', ACCENT_GREEN),
    ('', 'PPTX AI Assistant v0.1.0', LIGHT_GRAY),
    ('', 'Type /help for commands', SUBTLE_GRAY),
    ('', '', WHITE),
    ('You> ', '/help', ACCENT_BLUE),
    ('', '  /clear  会話履歴をクリア', LIGHT_GRAY),
    ('', '  /help   ヘルプを表示', LIGHT_GRAY),
    ('', '  /exit   終了', LIGHT_GRAY),
    ('', '  Ctrl+D  終了（代替）', LIGHT_GRAY),
    ('', '', WHITE),
    ('You> ', 'AI動向のプレゼンを作って', ACCENT_BLUE),
    ('', '  [generate_pptx] を実行中...', ACCENT_ORANGE),
]
for i, (prompt, text, color) in enumerate(cmds):
    y = Inches(2.55) + Inches(i * 0.35)
    if prompt:
        add_text(s, Inches(1.1), y, Inches(0.8), Inches(0.35),
                 prompt, font_size=14, color=ACCENT_GREEN, font_name='Consolas')
        add_text(s, Inches(1.9), y, Inches(5.5), Inches(0.35),
                 text, font_size=14, color=color, font_name='Consolas')
    else:
        add_text(s, Inches(1.1), y, Inches(6.5), Inches(0.35),
                 text, font_size=14, color=color, font_name='Consolas')

# UI features
add_text(s, Inches(9.0), Inches(2.0), Inches(4), Inches(0.5),
         'UI の特徴', font_size=22, color=WHITE, bold=True)
ui_feats = [
    ('Rich Panel', 'ツール呼び出しを色付きパネルで表示', ACCENT_BLUE),
    ('Markdown', 'AI応答をMarkdownでレンダリング', ACCENT_GREEN),
    ('History', '~/.cli_agent/history.txt に履歴保存', ACCENT_PURPLE),
    ('Truncation', '長い出力を自動的に省略表示', ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(ui_feats):
    y = Inches(2.8) + Inches(i * 1.1)
    add_shape(s, Inches(9.0), y, Inches(3.8), Inches(0.9),
              BG_CARD, border_color=color)
    add_text(s, Inches(9.3), y + Inches(0.05), Inches(3.3), Inches(0.4),
             title, font_size=15, color=color, bold=True, font_name='Consolas')
    add_text(s, Inches(9.3), y + Inches(0.4), Inches(3.3), Inches(0.4),
             desc, font_size=12, color=LIGHT_GRAY)
add_slide_number(s, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8: PPTX Generation Details
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(0.9), Inches(8), Inches(0.7),
         'PPTX 生成機能の詳細', font_size=36, color=WHITE, bold=True)
stypes = [
    ('content', 'コンテンツスライド',
     '見出し + 箇条書き\n標準的な情報伝達用\n4-6 項目推奨', ACCENT_BLUE),
    ('two_column', '2カラムスライド',
     '左右分割レイアウト\n比較・対比に最適\n各カラム独立制御', ACCENT_GREEN),
    ('title', 'タイトルスライド',
     '中央配置の大見出し\nサブタイトル対応\n44pt / 24pt', ACCENT_PURPLE),
]
for i, (tn, label, desc, color) in enumerate(stypes):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(2.0)
    add_shape(s, x, y, Inches(3.8), Inches(2.5), BG_CARD, border_color=color)
    add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.4),
             tn, font_size=18, color=color, bold=True, font_name='Consolas')
    add_text(s, x + Inches(0.3), y + Inches(0.65), Inches(3.3), Inches(0.35),
             label, font_size=16, color=WHITE, bold=True)
    add_text(s, x + Inches(0.3), y + Inches(1.1), Inches(3.3), Inches(1.2),
             desc, font_size=14, color=LIGHT_GRAY)

add_text(s, Inches(0.8), Inches(4.9), Inches(4), Inches(0.5),
         '生成仕様', font_size=22, color=WHITE, bold=True)
specs = [
    ('スライドサイズ', '10" x 7.5"（4:3）'),
    ('タイトル文字', '44pt Bold（中央配置）'),
    ('見出し文字', '32pt Bold'),
    ('本文文字', '18pt（行間 14pt）'),
    ('推奨スライド数', '6 〜 10 枚'),
    ('箇条書き', '自動 \u2022 付与（4-6項目）'),
]
for i, (key, val) in enumerate(specs):
    col, row = i % 2, i // 2
    x = Inches(0.8) + Inches(col * 6.0)
    y = Inches(5.5) + Inches(row * 0.55)
    add_text(s, x, y, Inches(2.2), Inches(0.45),
             key, font_size=15, color=ACCENT_BLUE, bold=True)
    add_text(s, x + Inches(2.2), y, Inches(3.5), Inches(0.45),
             val, font_size=15, color=LIGHT_GRAY, font_name='Consolas')
add_slide_number(s, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9: Configuration
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_GREEN)
add_text(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
         '設定 & ディレクトリ構成', font_size=36, color=WHITE, bold=True)

# Config card
add_shape(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0),
          BG_CARD, border_color=ACCENT_GREEN)
add_text(s, Inches(1.1), Inches(2.15), Inches(4), Inches(0.5),
         '設定項目', font_size=20, color=ACCENT_GREEN, bold=True)
cfgs = [
    ('ANTHROPIC_API_KEY', '.env ファイルで設定（必須）'),
    ('CLI_MODEL', 'モデル指定（デフォルト: claude-sonnet-4-20250514）'),
    ('MAX_TOKENS', '最大トークン数: 4096'),
    ('SYSTEM_PROMPT', '日本語のシステムプロンプト'),
    ('history_dir', '~/.cli_agent/（入力履歴保存）'),
]
for i, (key, desc) in enumerate(cfgs):
    y = Inches(2.8) + Inches(i * 0.8)
    add_text(s, Inches(1.1), y, Inches(5.0), Inches(0.35),
             key, font_size=14, color=ACCENT_BLUE, bold=True, font_name='Consolas')
    add_text(s, Inches(1.1), y + Inches(0.3), Inches(5.0), Inches(0.35),
             desc, font_size=13, color=LIGHT_GRAY)

# Directory card
add_shape(s, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0),
          BG_CARD, border_color=ACCENT_PURPLE)
add_text(s, Inches(7.1), Inches(2.15), Inches(4), Inches(0.5),
         'ディレクトリ構成', font_size=20, color=ACCENT_PURPLE, bold=True)
dirs = [
    'cli/',
    '  \u251c\u2500\u2500 app.py          # REPL ループ',
    '  \u251c\u2500\u2500 agent.py        # エージェントループ',
    '  \u251c\u2500\u2500 config.py       # 設定・プロンプト',
    '  \u251c\u2500\u2500 tools/',
    '  \u2502   \u251c\u2500\u2500 base.py     # ToolDefinition',
    '  \u2502   \u251c\u2500\u2500 file_tools.py',
    '  \u2502   \u251c\u2500\u2500 command_tool.py',
    '  \u2502   \u2514\u2500\u2500 pptx_tool.py',
    '  \u2514\u2500\u2500 ui/',
    '      \u251c\u2500\u2500 console.py  # Rich レンダリング',
    '      \u2514\u2500\u2500 input.py    # 入力管理',
    'scripts/              # ユーティリティ',
    '.claude/              # ドキュメント',
]
for i, line in enumerate(dirs):
    y = Inches(2.75) + Inches(i * 0.3)
    add_text(s, Inches(7.1), y, Inches(5.2), Inches(0.3),
             line, font_size=11, color=LIGHT_GRAY, font_name='Consolas')
add_slide_number(s, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10: Summary
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_shape(s, Inches(10), Inches(-1), Inches(5), Inches(5), BG_CARD)
add_accent_line(s, Inches(0.8), Inches(0.8), Inches(2.5), ACCENT_BLUE)
add_text(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.7),
         'まとめ', font_size=36, color=WHITE, bold=True)
summary = [
    ('自然言語 → PPTX',
     'Claude AI が指示を解釈し、構造化されたプレゼンテーションを自動生成',
     ACCENT_BLUE),
    ('6種の統合ツール',
     'ファイル操作・コマンド実行・PPTX生成を統合した包括的なツールセット',
     ACCENT_GREEN),
    ('自律エージェント',
     '最大20回のツール呼び出しチェーンで複雑なタスクを自律的に完遂',
     ACCENT_PURPLE),
    ('拡張可能な設計',
     'ツール追加・モデル切替・プロンプトカスタマイズが容易なモジュラー構成',
     ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(summary):
    y = Inches(2.0) + Inches(i * 1.3)
    add_shape(s, Inches(0.8), y, Inches(11.7), Inches(1.1),
              BG_CARD, border_color=color)
    add_shape(s, Inches(0.8), y, Pt(5), Inches(1.1), color)
    add_text(s, Inches(1.3), y + Inches(0.1), Inches(4), Inches(0.4),
             title, font_size=20, color=color, bold=True)
    add_text(s, Inches(1.3), y + Inches(0.55), Inches(10.5), Inches(0.45),
             desc, font_size=16, color=LIGHT_GRAY)
add_text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
         'python -m cli  |  v0.1.0  |  Python 3.8+  |  Claude AI Powered',
         font_size=16, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER,
         font_name='Consolas')
add_slide_number(s, 10, TOTAL_SLIDES)

# ========== Save ==========
output = 'pptx_ai_assistant_cli_spec.pptx'
prs.save(output)
print(f'Saved: {output}')
print(f'Slides: {len(prs.slides)}')
