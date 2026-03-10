"""PowerPoint generation tool using python-pptx — designed slides."""

import os
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from .base import ToolDefinition, ToolResult

# ──────────────────── Design Tokens ────────────────────
# Midnight Executive palette (default)
COLORS = {
    "primary": RGBColor(0x1E, 0x27, 0x61),   # Dark navy
    "secondary": RGBColor(0xCA, 0xDC, 0xFC),  # Ice blue
    "accent": RGBColor(0xF9, 0x61, 0x67),     # Coral
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_bg": RGBColor(0xF4, 0xF6, 0xFA),   # Light gray-blue
    "body_text": RGBColor(0x2D, 0x34, 0x36),   # Dark gray
    "muted": RGBColor(0x8E, 0x99, 0xA4),       # Muted gray
    "card_bg": RGBColor(0xFF, 0xFF, 0xFF),      # Card white
    "divider": RGBColor(0xE2, 0xE8, 0xF0),     # Light divider
}

FONT_HEADER = "Georgia"
FONT_BODY = "Calibri"

# Chart-specific colors (hex strings for chart API)
CHART_COLORS_HEX = ["1E2761", "F96167", "CADCFC", "2C5F2D", "F9E795", "065A82"]


def _add_accent_bar(slide, x, y, w, h, color=None):
    """Add a thin accent bar shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or COLORS["accent"]
    shape.line.fill.background()


def _add_bg_shape(slide, x, y, w, h, color):
    """Add a background rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_slide_bg(slide, color):
    """Set solid slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _build_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Add a beautifully designed dark title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLORS["primary"])

    # Left accent bar
    _add_accent_bar(slide, 0, 0, 0.12, 5.625)

    # Bottom decorative line
    _add_bg_shape(slide, 0, 4.8, 10, 0.06, COLORS["secondary"])

    # Small accent dot
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        Inches(0.8), Inches(1.8), Inches(0.15), Inches(0.15),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["accent"]
    shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.1), Inches(8.4), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONT_HEADER

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.8)
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLORS["secondary"]
        p2.font.name = FONT_BODY


def _build_content_slide(
    prs: Presentation,
    heading: str,
    bullets: list[str],
) -> None:
    """Add a styled content slide with heading and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLORS["light_bg"])

    # Top header bar
    _add_bg_shape(slide, 0, 0, 10, 1.0, COLORS["primary"])

    # Accent accent on header
    _add_accent_bar(slide, 0, 0.9, 2.5, 0.06)

    # Heading (on dark header)
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONT_HEADER

    # Content card background
    _add_bg_shape(slide, 0.4, 1.3, 9.2, 3.9, COLORS["card_bg"])

    # Bullet content
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(8.6), Inches(3.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[-1]
        # Clean bullet prefix
        clean = text.lstrip("•-* ")
        p.text = f"▸ {clean}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["body_text"]
        p.font.name = FONT_BODY
        p.space_after = Pt(12)


def _build_two_column_slide(
    prs: Presentation,
    heading: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    """Add a styled two-column slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLORS["light_bg"])

    # Top header bar
    _add_bg_shape(slide, 0, 0, 10, 1.0, COLORS["primary"])
    _add_accent_bar(slide, 0, 0.9, 2.5, 0.06)

    # Heading
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONT_HEADER

    # Left card
    _add_bg_shape(slide, 0.4, 1.3, 4.3, 3.9, COLORS["card_bg"])
    _add_accent_bar(slide, 0.4, 1.3, 0.08, 3.9, COLORS["accent"])

    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(3.8), Inches(3.5)
    )
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.font.name = FONT_HEADER
    p.space_after = Pt(8)
    for text in left_bullets:
        p = tf_l.add_paragraph()
        clean = text.lstrip("•-* ")
        p.text = f"▸ {clean}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["body_text"]
        p.font.name = FONT_BODY
        p.space_after = Pt(8)

    # Right card
    _add_bg_shape(slide, 5.3, 1.3, 4.3, 3.9, COLORS["card_bg"])
    _add_accent_bar(slide, 5.3, 1.3, 0.08, 3.9, COLORS["secondary"])

    right_box = slide.shapes.add_textbox(
        Inches(5.6), Inches(1.5), Inches(3.8), Inches(3.5)
    )
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.font.name = FONT_HEADER
    p.space_after = Pt(8)
    for text in right_bullets:
        p = tf_r.add_paragraph()
        clean = text.lstrip("•-* ")
        p.text = f"▸ {clean}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["body_text"]
        p.font.name = FONT_BODY
        p.space_after = Pt(8)


def _build_chart_slide(
    prs: Presentation,
    heading: str,
    chart_type: str,
    categories: list[str],
    series: list[dict],
) -> None:
    """Add a slide with a styled chart (bar, pie, or line)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLORS["light_bg"])

    # Header bar
    _add_bg_shape(slide, 0, 0, 10, 1.0, COLORS["primary"])
    _add_accent_bar(slide, 0, 0.9, 2.5, 0.06)

    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONT_HEADER

    # Chart type mapping
    ct_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "pie": XL_CHART_TYPE.PIE,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    xl_type = ct_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name", "Data"), s.get("values", []))

    # Card background for chart
    _add_bg_shape(slide, 0.4, 1.2, 9.2, 4.1, COLORS["card_bg"])

    chart_frame = slide.shapes.add_chart(
        xl_type, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.7), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # Style chart colors
    plot = chart.plots[0]
    for idx, s in enumerate(plot.series):
        color_hex = CHART_COLORS_HEX[idx % len(CHART_COLORS_HEX)]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)


def _build_stat_slide(
    prs: Presentation,
    heading: str,
    stats: list[dict],
) -> None:
    """Add a slide with large stat callouts (big numbers)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLORS["primary"])

    # Heading
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONT_HEADER

    # Bottom accent line
    _add_bg_shape(slide, 0, 4.8, 10, 0.06, COLORS["secondary"])

    # Stat cards (up to 4)
    display_stats = stats[:4]
    n = len(display_stats)
    card_w = min(2.0, (9.0 - 0.3 * (n - 1)) / n)
    total_w = card_w * n + 0.3 * (n - 1)
    start_x = (10 - total_w) / 2

    accent_colors = [COLORS["accent"], COLORS["secondary"],
                     RGBColor(0x2C, 0x5F, 0x2D), RGBColor(0xF9, 0xE7, 0x95)]

    for i, stat in enumerate(display_stats):
        x = start_x + i * (card_w + 0.3)
        # Card bg
        _add_bg_shape(slide, x, 1.3, card_w, 3.0, RGBColor(0x2A, 0x31, 0x6E))
        # Top accent bar
        _add_accent_bar(slide, x, 1.3, card_w, 0.06, accent_colors[i % len(accent_colors)])

        # Big number
        num_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(1.7), Inches(card_w - 0.2), Inches(1.2)
        )
        tf_n = num_box.text_frame
        tf_n.word_wrap = True
        pn = tf_n.paragraphs[0]
        pn.text = str(stat.get("value", "0"))
        pn.font.size = Pt(44)
        pn.font.bold = True
        pn.font.color.rgb = accent_colors[i % len(accent_colors)]
        pn.font.name = FONT_HEADER
        pn.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(3.0), Inches(card_w - 0.2), Inches(1.0)
        )
        tf_l = lbl_box.text_frame
        tf_l.word_wrap = True
        pl = tf_l.paragraphs[0]
        pl.text = stat.get("label", "")
        pl.font.size = Pt(14)
        pl.font.color.rgb = COLORS["secondary"]
        pl.font.name = FONT_BODY
        pl.alignment = PP_ALIGN.CENTER


def _build_timeline_slide(
    prs: Presentation,
    heading: str,
    steps: list[dict],
) -> None:
    """Add a horizontal timeline / process flow slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLORS["light_bg"])

    # Header bar
    _add_bg_shape(slide, 0, 0, 10, 1.0, COLORS["primary"])
    _add_accent_bar(slide, 0, 0.9, 2.5, 0.06)

    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.font.name = FONT_HEADER

    display_steps = steps[:5]
    n = len(display_steps)
    step_w = min(1.8, (9.0 - 0.2 * (n - 1)) / n)
    total_w = step_w * n + 0.2 * (n - 1)
    start_x = (10 - total_w) / 2

    # Horizontal connector line
    connector_y = 2.6
    _add_bg_shape(slide, start_x, connector_y, total_w, 0.04, COLORS["divider"])

    for i, step in enumerate(display_steps):
        x = start_x + i * (step_w + 0.2)

        # Circle number
        circ = slide.shapes.add_shape(
            9,  # MSO_SHAPE.OVAL
            Inches(x + step_w / 2 - 0.2), Inches(connector_y - 0.18),
            Inches(0.4), Inches(0.4),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = COLORS["accent"] if i % 2 == 0 else COLORS["primary"]
        circ.line.fill.background()

        # Step number in circle
        circ.text_frame.paragraphs[0].text = str(i + 1)
        circ.text_frame.paragraphs[0].font.size = Pt(12)
        circ.text_frame.paragraphs[0].font.bold = True
        circ.text_frame.paragraphs[0].font.color.rgb = COLORS["white"]
        circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Step title
        t_box = slide.shapes.add_textbox(
            Inches(x), Inches(3.1), Inches(step_w), Inches(0.5)
        )
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        pt = tf_t.paragraphs[0]
        pt.text = step.get("title", "")
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = COLORS["primary"]
        pt.font.name = FONT_HEADER
        pt.alignment = PP_ALIGN.CENTER

        # Step description
        d_box = slide.shapes.add_textbox(
            Inches(x), Inches(3.6), Inches(step_w), Inches(1.5)
        )
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        pd = tf_d.paragraphs[0]
        pd.text = step.get("description", "")
        pd.font.size = Pt(11)
        pd.font.color.rgb = COLORS["body_text"]
        pd.font.name = FONT_BODY
        pd.alignment = PP_ALIGN.CENTER


async def _generate_pptx(
    filename: str,
    title: str,
    subtitle: str = "",
    slides: list[dict[str, Any]] | None = None,
) -> ToolResult:
    """Generate a professionally designed PowerPoint file."""
    if slides is None:
        slides = []

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    _build_title_slide(prs, title, subtitle)

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        if slide_type == "two_column":
            _build_two_column_slide(
                prs,
                heading=slide_data.get("heading", ""),
                left_title=slide_data.get("left_title", ""),
                left_bullets=slide_data.get("left_bullets", []),
                right_title=slide_data.get("right_title", ""),
                right_bullets=slide_data.get("right_bullets", []),
            )
        elif slide_type == "chart":
            _build_chart_slide(
                prs,
                heading=slide_data.get("heading", ""),
                chart_type=slide_data.get("chart_type", "bar"),
                categories=slide_data.get("categories", []),
                series=slide_data.get("series", []),
            )
        elif slide_type == "stat":
            _build_stat_slide(
                prs,
                heading=slide_data.get("heading", ""),
                stats=slide_data.get("stats", []),
            )
        elif slide_type == "timeline":
            _build_timeline_slide(
                prs,
                heading=slide_data.get("heading", ""),
                steps=slide_data.get("steps", []),
            )
        else:
            _build_content_slide(
                prs,
                heading=slide_data.get("heading", ""),
                bullets=slide_data.get("bullets", []),
            )

    if not filename.endswith(".pptx"):
        filename = f"{filename}.pptx"
    project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    output_dir = os.path.join(project_root, "output")
    os.makedirs(output_dir, exist_ok=True)
    abs_path = os.path.join(output_dir, os.path.basename(filename))
    prs.save(abs_path)
    slide_count = len(prs.slides)
    return f"PowerPoint generated: {abs_path} ({slide_count} slides)"


generate_pptx_tool = ToolDefinition(
    name="generate_pptx",
    description=(
        "Generate a professionally designed PowerPoint (.pptx) file. "
        "Slide types: 'content' (bullet points), 'two_column' (comparison), "
        "'chart' (bar/pie/line/doughnut), 'stat' (big number callouts), "
        "'timeline' (process flow). Use 'chart' for data visualization, "
        "'stat' for key metrics, and 'timeline' for step-by-step processes. "
        "Mix different types for engaging presentations."
    ),
    input_schema={
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Output filename (e.g., 'presentation.pptx')",
            },
            "title": {
                "type": "string",
                "description": "Presentation title (shown on dark title slide)",
            },
            "subtitle": {
                "type": "string",
                "description": "Subtitle on title slide",
                "default": "",
            },
            "slides": {
                "type": "array",
                "description": "List of slide definitions",
                "items": {
                    "type": "object",
                    "properties": {
                        "type": {
                            "type": "string",
                            "enum": ["content", "two_column", "chart", "stat", "timeline"],
                            "description": "Slide type: content, two_column, chart, stat, timeline",
                            "default": "content",
                        },
                        "heading": {"type": "string", "description": "Slide heading"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "For 'content': bullet points",
                        },
                        "left_title": {"type": "string"},
                        "left_bullets": {"type": "array", "items": {"type": "string"}},
                        "right_title": {"type": "string"},
                        "right_bullets": {"type": "array", "items": {"type": "string"}},
                        "chart_type": {
                            "type": "string",
                            "enum": ["bar", "pie", "line", "doughnut"],
                            "description": "For 'chart': chart visualization type",
                        },
                        "categories": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "For 'chart': category labels (x-axis)",
                        },
                        "series": {
                            "type": "array",
                            "description": "For 'chart': data series [{name, values}]",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": "string"},
                                    "values": {"type": "array", "items": {"type": "number"}},
                                },
                            },
                        },
                        "stats": {
                            "type": "array",
                            "description": "For 'stat': key metrics [{value, label}]",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "value": {"type": "string", "description": "Big number (e.g. '87%', '$2.4M')"},
                                    "label": {"type": "string", "description": "Description below number"},
                                },
                            },
                        },
                        "steps": {
                            "type": "array",
                            "description": "For 'timeline': process steps [{title, description}]",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "description": {"type": "string"},
                                },
                            },
                        },
                    },
                    "required": ["heading"],
                },
            },
        },
        "required": ["filename", "title"],
    },
    handler=_generate_pptx,
)
