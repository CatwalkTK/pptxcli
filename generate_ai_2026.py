#!/usr/bin/env python3
"""
2026年のAI事情 - スライド生成スクリプト
animations.md に基づき Fade トランジションを適用
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ========== スライド1: タイトル ==========
title_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "2026年のAI事情"
p.font.size = Pt(44)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

sub_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.text = "自律型エージェント・本格実装・新たな働き方"
p2.font.size = Pt(24)
p2.alignment = PP_ALIGN.CENTER

# ========== スライド2: 2026年 主要トレンド ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "2026年 AI 主要トレンド"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
items = [
    ("1. エージェント型AIへの進化", "受動的チャットから自律的に目標達成するAIへ"),
    ("2. 企業の本格実装フェーズ", "PoCを終え、ROIを実感できる段階へ"),
    ("3. インフラ効率化", "ASIC・チップレット・量子支援型オプティマイザー"),
    ("4. 生成動画の本格普及", "映像制作業界の構造変化が進行"),
    ("5. プライバシー重視のAI", "デバイス上でのデータ処理が注目"),
]
for i, (title_text, desc) in enumerate(items):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = title_text
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = f"  -> {desc}"
    p2.font.size = Pt(14)
    p2.level = 1
    p2.space_after = Pt(12)

# ========== スライド3: 自律型AIエージェント ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "自律型AIエージェントの台頭"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "2025年76億ドル -> 2026年約111億ドルへ急成長（AIエージェント市場）",
    "単なる自動化ツールから「自律的行為主体」へ進化",
    "セルフリファイン：AIが自ら評価・改善を繰り返し品質向上",
    "マルチエージェント、API統合型、OS統合型が並行展開",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド4: 企業実装 ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "企業の本格実装フェーズ"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "96%の組織がAIエージェント活用の拡大を計画",
    "67%の経営者が「12ヶ月以内に役割を変革」と同意",
    "グローバル企業は全社エージェント化を推進",
    "PoCから本格導入・ROI可視化の段階へ",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド5: 働き方の変化 ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "働き方の変革"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "人間の役割：「作業実行者」->「目標設定者・最終判断者」へ",
    "複数エージェントを指揮する「オーケストラ指揮者」的な役割",
    "AIエージェントとの協働を支える新職種の価値確立",
    "各エージェントにデジタルID、アクセス権限管理の標準化",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド6: 技術・インフラ ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "インフラ・技術トレンド"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "ASICベースのアクセラレーター、チップレット設計の成熟",
    "量子支援型オプティマイザー（IBM: 2026年が量子超越の年と予測）",
    "合成データ活用：プライバシーを守りながらAI訓練",
    "主要フレームワーク：LangGraph、CrewAI、AutoGPT",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド7: 課題とガバナンス ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "課題とガバナンス"
p.font.size = Pt(32)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
tf = content_box.text_frame
points = [
    "安全設計、説明責任、人間中心設計が喫緊の課題",
    "著作権問題の激化：クリエイターとAI企業の対立",
    "セキュリティは設計段階から組み込みが必須",
    "自律性の向上とガードレールのバランス",
]
for i, text in enumerate(points):
    if i > 0:
        tf.add_paragraph()
    p = tf.paragraphs[-1]
    p.text = f"• {text}"
    p.font.size = Pt(18)
    p.space_after = Pt(14)

# ========== スライド8: まとめ ==========
content_slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "まとめ"
p.font.size = Pt(36)
p.font.bold = True

content_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
tf = content_box.text_frame
p = tf.paragraphs[0]
p.text = "2026年：AIエージェントが「当たり前」になる年"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(20)

for text in [
    "• 自律型エージェントが企業の標準インフラに",
    "• 人間の役割は目標設定と最終判断へシフト",
    "• インフラ効率化で量子・ASICが本格化",
    "• ガバナンスと安全設計が差別化要因に",
]:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.space_after = Pt(10)

# 一時ファイルとして保存（output/ディレクトリへ）
output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
os.makedirs(output_dir, exist_ok=True)
static_path = os.path.join(output_dir, "ai_2026_static.pptx")
prs.save(static_path)
print(f"静的PPTXを生成: {static_path}")
